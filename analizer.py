import os
import json
from docx import Document
from pypdf import PdfReader
from pathlib import Path
from pptx import Presentation

from image import get_image_annotation
from video import get_video_annotation


class Analizer:
    def __init__(self, client, model):
        self.client = client
        self.model = model

    
    def getText(self, path_to_dir:str):
        text = ""
        folder = Path(path_to_dir)
        for file in folder.rglob("*"):
            if not file.is_file():
                continue

            ext = file.suffix.lower()
            print("extracting text: ", ext)
            
            if ext in [".md", ".txt"]:
                with open(file, "r", encoding="utf-8") as f:
                    text += file.name + ":\n\n" + f.read() + "\n\n"

            if ext == ".pdf":
                reader = PdfReader(file)
                text += file.name + ":\n\n" + "\n".join(page.extract_text() or "" for page in reader.pages) + "\n\n"

            if ext in [".doc", ".docx"]:
                doc_text = ""
                doc = Document(file)
                # Извлекаем элементы документа в порядке их появления
                for element in doc.element.body:
                    if element.tag.endswith('p'):
                        # Это параграф
                        for paragraph in doc.paragraphs:
                            if paragraph._element == element:
                                doc_text += paragraph.text + "\n"
                                break
                    elif element.tag.endswith('tbl'):
                        # Это таблица
                        for table in doc.tables:
                            if table._element == element:
                                doc_text += "\n[Таблица]\n"
                                for row in table.rows:
                                    row_text = " | ".join(cell.text for cell in row.cells)
                                    doc_text += row_text + "\n"
                                doc_text += "\n"
                                break
                text += file.name + ":\n\n" + doc_text + "\n\n"

            if ext == ".ipynb":
                text += file.name + ":\n\n" + extract_ipynb_text(file.read_text(encoding="utf-8")) + "\n\n"

            if ext == ".pptx":
                text += file.name + ":\n\n" + extract_text_from_pptx(file)
        return text
    

    def getMedia(self, path_to_dir:str):
        text = ""
        folder = Path(path_to_dir)
        for file in folder.rglob("*"):
            if not file.is_file():
                continue

            ext = file.suffix.lower()

            if ext in [".jpeg", ".png", ".webp"]:
                text += file.name + ":\n\n" + get_image_annotation(file, self.client)

            if ext in [".mp4", ".avi", ".mov", ".mkv", ".wmv", ".webm"]:
                text += file.name + ":\n\n" + get_video_annotation(file, self.client)
        return text


def extract_ipynb_text(text):
    data = json.loads(text)
    parts = []

    for cell in data.get("cells", []):
        src = cell.get("source", [])
        if isinstance(src, list):
            text = "".join(src)
        else:
            text = src
        if text.strip():
            parts.append(text)

    return "\n\n".join(parts)

from pptx import Presentation


def extract_text_from_pptx(file_path: str) -> str:
    """
    Извлекает весь текст из презентации PowerPoint.
    
    :param file_path: путь до .pptx файла
    :return: строка со всем текстом из презентации
    """
    prs = Presentation(file_path)
    result = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        result.append(f"## Слайд {slide_num}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        result.append(text)

    return "\n".join(result)